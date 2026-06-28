import datetime
import json
import os
import tempfile
from dataclasses import dataclass

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt


PALETTE = {
    "primary": RGBColor(0x1A, 0x1A, 0x2E),  # deep navy
    "accent": RGBColor(0x16, 0x21, 0x3E),  # dark blue
    "highlight": RGBColor(0xE9, 0x4F, 0x37),  # red-orange
    "light": RGBColor(0xF5, 0xF5, 0xF5),  # off-white
    "mid": RGBColor(0x9E, 0x9E, 0x9E),  # grey
}

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


@dataclass
class SlideResult:
    file_path: str
    title: str


def _new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    for layout in prs.slide_master.slide_layouts:
        if "blank" in layout.name.lower():
            return prs.slides.add_slide(layout)

    return prs.slides.add_slide(prs.slide_layouts[6])


def _add_textbox(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    text: str,
    font_size: int,
    bold: bool = False,
    color: RGBColor = None,
    align=PP_ALIGN.LEFT,
    italic: bool = False,
) -> None:
    textbox = slide.shapes.add_textbox(
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align

    font = paragraph.font
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    if color is not None:
        font.color.rgb = color


def _fill_slide_background(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _build_title_slide(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["primary"])

    _add_textbox(
        slide,
        left=0,
        top=0.75,
        width=13.33,
        height=3.0,
        text=slide_data["title"],
        font_size=54,
        bold=True,
        color=PALETTE["light"],
        align=PP_ALIGN.CENTER,
    )

    subtitle = slide_data.get("content", {}).get("subtitle", "")
    if subtitle:
        _add_textbox(
            slide,
            left=1.0,
            top=4.0,
            width=11.33,
            height=1.0,
            text=subtitle,
            font_size=28,
            color=PALETTE["mid"],
            align=PP_ALIGN.CENTER,
        )

    return slide


def _build_title_and_content(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["light"])

    _add_textbox(
        slide,
        left=0.6,
        top=0.4,
        width=12,
        height=1.2,
        text=slide_data["title"],
        font_size=36,
        bold=True,
        color=PALETTE["primary"],
    )

    rule = slide.shapes.add_shape(
        1,
        0,
        Inches(1.6),
        SLIDE_W,
        Inches(0.04),
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = PALETTE["highlight"]
    rule.line.fill.background()

    body = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(1.8),
        Inches(12),
        Inches(5.2),
    )
    tf = body.text_frame
    tf.word_wrap = True
    bullets = slide_data.get("content", {}).get("bullets", [])

    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.font.size = Pt(24)
        paragraph.font.color.rgb = PALETTE["accent"]

    return slide


def _build_big_statement(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["primary"])

    statement = slide_data.get("content", {}).get("statement", slide_data["title"])
    _add_textbox(
        slide,
        left=1.0,
        top=2.0,
        width=11.33,
        height=3.5,
        text=statement,
        font_size=54,
        bold=True,
        color=PALETTE["light"],
        align=PP_ALIGN.CENTER,
    )

    bar = slide.shapes.add_shape(
        1,
        Inches(5.9),
        Inches(5.8),
        Inches(1.5),
        Inches(0.08),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PALETTE["highlight"]
    bar.line.fill.background()

    return slide


def _build_section_header(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["accent"])

    bar = slide.shapes.add_shape(
        1,
        0,
        0,
        Inches(0.15),
        Inches(7.5),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PALETTE["highlight"]
    bar.line.fill.background()

    heading = slide_data.get("content", {}).get("heading", slide_data["title"])
    _add_textbox(
        slide,
        left=0.8,
        top=2.5,
        width=12,
        height=1.8,
        text=heading,
        font_size=48,
        bold=True,
        color=PALETTE["light"],
        align=PP_ALIGN.LEFT,
    )

    subheading = slide_data.get("content", {}).get("subheading", "")
    _add_textbox(
        slide,
        left=0.8,
        top=4.4,
        width=10,
        height=1.0,
        text=subheading,
        font_size=24,
        color=PALETTE["mid"],
        align=PP_ALIGN.LEFT,
    )

    return slide


def _build_closing(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["primary"])

    bar = slide.shapes.add_shape(
        1,
        0,
        Inches(3.4),
        SLIDE_W,
        Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PALETTE["highlight"]
    bar.line.fill.background()

    headline = slide_data.get("content", {}).get("headline", slide_data["title"])
    _add_textbox(
        slide,
        left=1.16,
        top=1.5,
        width=11,
        height=1.8,
        text=headline,
        font_size=48,
        bold=True,
        color=PALETTE["light"],
        align=PP_ALIGN.CENTER,
    )

    cta = slide_data.get("content", {}).get("cta", "")
    _add_textbox(
        slide,
        left=2.16,
        top=4.0,
        width=9,
        height=1.2,
        text=cta,
        font_size=28,
        bold=True,
        color=PALETTE["highlight"],
        align=PP_ALIGN.CENTER,
    )

    return slide


def _build_two_column(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["light"])
    content = slide_data.get("content", {})

    _add_textbox(
        slide,
        left=0.6,
        top=0.4,
        width=12,
        height=1.1,
        text=slide_data["title"],
        font_size=34,
        bold=True,
        color=PALETTE["primary"],
    )

    divider = slide.shapes.add_shape(
        1,
        Inches(6.55),
        Inches(1.7),
        Inches(0.05),
        Inches(5.3),
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = PALETTE["mid"]
    divider.line.fill.background()

    _add_textbox(
        slide,
        left=0.5,
        top=1.6,
        width=5.8,
        height=0.7,
        text=content["left_header"],
        font_size=22,
        bold=True,
        color=PALETTE["highlight"],
    )

    left_box = slide.shapes.add_textbox(
        Inches(0.5),
        Inches(2.4),
        Inches(5.8),
        Inches(4.5),
    )
    left_tf = left_box.text_frame
    left_tf.word_wrap = True
    for index, bullet in enumerate(content["left_bullets"]):
        paragraph = left_tf.paragraphs[0] if index == 0 else left_tf.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = PALETTE["accent"]

    _add_textbox(
        slide,
        left=6.8,
        top=1.6,
        width=5.8,
        height=0.7,
        text=content["right_header"],
        font_size=22,
        bold=True,
        color=PALETTE["primary"],
    )

    right_box = slide.shapes.add_textbox(
        Inches(6.8),
        Inches(2.4),
        Inches(5.8),
        Inches(4.5),
    )
    right_tf = right_box.text_frame
    right_tf.word_wrap = True
    for index, bullet in enumerate(content["right_bullets"]):
        paragraph = right_tf.paragraphs[0] if index == 0 else right_tf.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = PALETTE["accent"]

    return slide


def _build_quote(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["accent"])
    content = slide_data.get("content", {})

    _add_textbox(
        slide,
        left=0.5,
        top=0.3,
        width=3,
        height=2.5,
        text="“",
        font_size=120,
        bold=True,
        color=PALETTE["highlight"],
    )

    _add_textbox(
        slide,
        left=1.2,
        top=1.5,
        width=10.8,
        height=3.5,
        text=content["quote"],
        font_size=32,
        color=PALETTE["light"],
        align=PP_ALIGN.CENTER,
        italic=True,
    )

    _add_textbox(
        slide,
        left=1.2,
        top=5.2,
        width=10.8,
        height=0.8,
        text="— " + content.get("attribution", ""),
        font_size=20,
        color=PALETTE["mid"],
        align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide,
        left=12.5,
        top=4.5,
        width=2,
        height=1.5,
        text="”",
        font_size=120,
        bold=True,
        color=PALETTE["highlight"],
    )

    return slide


def _build_image_and_text(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["light"])
    content = slide_data.get("content", {})

    _add_textbox(
        slide,
        left=0.5,
        top=0.3,
        width=12.3,
        height=1.0,
        text=slide_data["title"],
        font_size=32,
        bold=True,
        color=PALETTE["primary"],
    )

    placeholder = slide.shapes.add_shape(
        1,
        Inches(0.5),
        Inches(1.5),
        Inches(5.8),
        Inches(5.5),
    )
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = PALETTE["mid"]
    placeholder.line.fill.background()

    _add_textbox(
        slide,
        left=0.5,
        top=3.8,
        width=5.8,
        height=0.7,
        text="[Image]",
        font_size=24,
        color=PALETTE["light"],
        align=PP_ALIGN.CENTER,
    )

    _add_textbox(
        slide,
        left=0.5,
        top=5.1,
        width=5.8,
        height=0.7,
        text=content.get("caption", ""),
        font_size=16,
        color=PALETTE["mid"],
        italic=True,
    )

    bullets_box = slide.shapes.add_textbox(
        Inches(6.8),
        Inches(1.5),
        Inches(6.0),
        Inches(5.5),
    )
    bullets_tf = bullets_box.text_frame
    bullets_tf.word_wrap = True
    for index, bullet in enumerate(content.get("bullets", [])):
        paragraph = bullets_tf.paragraphs[0] if index == 0 else bullets_tf.add_paragraph()
        paragraph.text = f"• {bullet}"
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = PALETTE["accent"]

    return slide


def _build_timeline(prs: Presentation, slide_data: dict):
    slide = _blank_slide(prs)
    _fill_slide_background(slide, PALETTE["primary"])

    _add_textbox(
        slide,
        left=0.6,
        top=0.3,
        width=12,
        height=1.0,
        text=slide_data["title"],
        font_size=34,
        bold=True,
        color=PALETTE["light"],
    )

    spine = slide.shapes.add_shape(
        1,
        Inches(6.6),
        Inches(1.4),
        Inches(0.1),
        Inches(5.7),
    )
    spine.fill.solid()
    spine.fill.fore_color.rgb = PALETTE["highlight"]
    spine.line.fill.background()

    events = slide_data.get("content", {}).get("events", [])[:6]

    for index, event in enumerate(events):
        top_inches = 1.5 + index * (5.5 / max(len(events) - 1, 1))
        top_inches = min(max(top_inches, 1.5), 6.8)
        node_top = Inches(top_inches)

        node = slide.shapes.add_shape(
            MSO_SHAPE_TYPE.OVAL,
            Inches(6.53),
            node_top - Inches(0.125),
            Inches(0.25),
            Inches(0.25),
        )
        node.fill.solid()
        node.fill.fore_color.rgb = PALETTE["highlight"]
        node.line.fill.background()

        is_left = index % 2 == 0
        text_left = 0.4 if is_left else 7.1
        align = PP_ALIGN.RIGHT if is_left else PP_ALIGN.LEFT

        _add_textbox(
            slide,
            left=text_left,
            top=top_inches - 0.3,
            width=5.8,
            height=0.5,
            text=event["label"],
            font_size=18,
            bold=True,
            color=PALETTE["highlight"],
            align=align,
        )

        _add_textbox(
            slide,
            left=text_left,
            top=top_inches + 0.25,
            width=5.8,
            height=0.5,
            text=event["description"],
            font_size=15,
            color=PALETTE["mid"],
            align=align,
        )

        connector = slide.shapes.add_shape(
            1,
            Inches(6.3 if is_left else 6.7),
            node_top - Inches(0.01),
            Inches(0.3),
            Inches(0.02),
        )
        connector.fill.solid()
        connector.fill.fore_color.rgb = PALETTE["highlight"]
        connector.line.fill.background()

    return slide


_LAYOUT_BUILDERS = {
    "title_slide": _build_title_slide,
    "title_and_content": _build_title_and_content,
    "big_statement": _build_big_statement,
    "section_header": _build_section_header,
    "closing": _build_closing,
    "two_column": _build_two_column,
    "quote": _build_quote,
    "image_and_text": _build_image_and_text,
    "timeline": _build_timeline,
    "data_table": None,  # reserved — falls back to title_and_content
}


def _inject_speaker_notes(slide, notes: list[str]) -> None:
    tf = slide.notes_slide.notes_text_frame
    tf.text = ""

    for index, note in enumerate(notes):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = f"• {note}"


def _safe_filename(title: str) -> str:
    stem = title.lower().replace(" ", "_")
    stem = "".join(char for char in stem if char.isalnum() or char == "_")
    stem = stem[:50] or "presentation"
    return f"{stem}.pptx"


def build_presentation(outline: dict, repo_source: str = "") -> SlideResult:
    output_dir = os.environ.get("SLIDES_OUTPUT_DIR") or tempfile.mkdtemp()
    prs = _new_presentation()

    _build_title_slide(
        prs,
        {"title": outline["title"], "content": {"subtitle": ""}},
    )

    for slide_data in outline["slides"]:
        layout = slide_data.get("layout", "title_and_content")
        builder = _LAYOUT_BUILDERS.get(layout)
        if builder is None:
            builder = _build_title_and_content

        slide = builder(prs, slide_data)
        _inject_speaker_notes(slide, slide_data.get("speaker_notes", []))

    filename = _safe_filename(outline["title"])
    full_output_path = os.path.join(output_dir, filename)
    prs.save(full_output_path)

    if repo_source:
        meta_path = os.path.splitext(full_output_path)[0] + ".meta.json"
        sidecar_path = meta_path
        meta = {
            "repo_source": repo_source,
            "presentation_title": outline["title"],
            "generated_at": datetime.datetime.utcnow().isoformat() + "Z",
            "slide_count": len(outline.get("slides", [])),
            "layout_summary": {
                layout: sum(1 for s in outline.get("slides", []) if s.get("layout") == layout)
                for layout in set(s.get("layout", "unknown") for s in outline.get("slides", []))
            },
        }
        with open(sidecar_path, "w", encoding="utf-8") as f:
            json.dump(meta, f, indent=2)

    return SlideResult(file_path=full_output_path, title=outline["title"])
